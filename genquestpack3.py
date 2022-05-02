# -*- coding: utf-8 -*-

# Import the modules    
import requests
import urllib3
from urllib.parse import urlparse
from urllib.parse import urljoin
from bs4 import BeautifulSoup
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from unidecode import unidecode
import io
import sys
import os

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

def download_image_old(url):
	urlo = urlparse(url)
	image_name = 'qimg.jpg'

    # fix url
	# TODO: urljoin
	if url[:1] == '/':
		url = urlo.scheme + '://' + urlo.netloc + url

	r = requests.get(url, stream=True, verify=False)
	with open(image_name, 'wb') as f:
		for chunk in r.iter_content():
			f.write(chunk)

# https://pythonprogramming.altervista.org/inserting-an-image-in-powerpoint-with-python/
def add_image(slide, placeholder_id, image_url):
	#for shape in slide.placeholders:
	#	print('%d %s' % (shape.placeholder_format.idx, shape.name))
	
	# to abs
	#if image_url[:1] == '/':
	#	image_url = urljoin()

	# download file
	r = requests.get(image_url)
	if r.status_code != 200:
		assert False, 'Status code error: {}.'.format(r.status_code)

	placeholder = slide.placeholders[placeholder_id]
 
    # Calculate the image size of the image
	im = Image.open(io.BytesIO(r.content))
	width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
	placeholder.height = height
	placeholder.width = width
 
	# Insert the picture
	placeholder = placeholder.insert_picture(io.BytesIO(r.content))
 
    # Calculate ratios and compare
	image_ratio = width / height
	placeholder_ratio = placeholder.width / placeholder.height
	ratio_difference = placeholder_ratio - image_ratio
 
	# Placeholder width too wide:
	if ratio_difference > 0:
		difference_on_each_side = ratio_difference / 2
		placeholder.crop_left = -difference_on_each_side
		placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
	else:
		difference_on_each_side = -ratio_difference / 2
		placeholder.crop_bottom = -difference_on_each_side
		placeholder.crop_top = -difference_on_each_side

def generate_pptx(url, file_stream = None):
	print('<- ', url)
	r = requests.get(url, verify=False)
	page = r.content.decode('utf-8')
	soup = BeautifulSoup(page, 'lxml')
	urlo = urlparse(url)
	if file_stream == None:
		prs = Presentation('braindo-tmpl.pptx')
	else:
		prs = Presentation(file_stream)

	if urlo.netloc == 'kand.info':
		# pack title
		pack_title = soup.find('h2', 'with-tabs').text

		# regexp 
		# TODO: <div class='razdatka_header'>Раздаточный материал</div> 
		question = 'Вопрос '
		answer = 'Ответ:'
		comment = 'Комментарий:'
		alt_answer = 'Зачёт:'
		questionimg = re.compile(r'/sites/default/files/.+')

		qsplit = '<br/><hr/><br/>'
		psplit = 'p'

	elif urlo.netloc == 'db.chgk.info':
		# pack title
		pack_title = soup.title.text

		# regexp 
		question = 'Вопрос '
		answer = 'Ответ:'
		comment = 'Комментарий:'
		alt_answer = 'Зачёт:'
		questionimg = re.compile(r'http://db.chgk.info/images/db/\d+.jpg')

		qsplit = '<div class="question" id=".+">'
		psplit = 'p'

	else:
		print ('site is not supported')
		exit(0)

	# Title Slide
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = pack_title
	subtitle.text = url

	print(pack_title)

	text_slide_layout = prs.slide_layouts[1]
	img_slide_layout = prs.slide_layouts[2]

	qs = re.split(qsplit, page) #TODO: SoupStrainer
	i = 1

	for q in qs:
		soup = BeautifulSoup(q, 'lxml') 

		for qp in soup(psplit): 
			
			txt = qp.get_text(strip=True)
			# Вопрос 1:
			m = txt.startswith(question)
			if m:
				print('Parsing question %d/%d' % (i, len(qs) - 1))
				i = i + 1

				if txt.partition(':')[2] != '':
					slide = prs.slides.add_slide(text_slide_layout)
					title = slide.shapes.title
					text = slide.placeholders[1]

					title.text = 'Вопрос'
					text.text = txt.partition(':')[2]
					if text.text == ' ':
						text.text = 'q is not found'
				# с раздаткой
				else:
					slide = prs.slides.add_slide(img_slide_layout)
					title = slide.shapes.title
					text = slide.placeholders[1]
					handout = slide.placeholders[13]

					title.text = 'Вопрос'
					text.text = qp.nextSibling.nextSibling.nextSibling.nextSibling
					if text.text == ' ':
						text.text = 'q is not found'

					# Вопросы - картинки
					#for img in soup.find_all('img', src = questionimg):
						#download_image(img['src'])				
						#pic = slide.shapes.add_picture('qimg.jpg', text.left, text.top)
						#text.top += pic.top
						#os.remove('qimg.jpg')
						# TODO: add_image
					if len(list(qp.nextSibling.children)) == 7 and soup.find('img', src = questionimg) != None:
						add_image(slide, 13, soup.find('img', src = questionimg)['src'])
					else:
						handout.text = '' #list(qp.nextSibling.children)[4]

			# Ответ:
			m = txt.startswith(answer)
			if m:
				slide = prs.slides.add_slide(text_slide_layout)
				title = slide.shapes.title
				text = slide.placeholders[1]

				title.text = 'Ответ'
				text.text = txt.partition(':')[2]

			# Зачёт: 
			m = txt.startswith(alt_answer)
			if m:
				# обычно идет после ответа поэтому отдельного слайда не требуется
				text.text += '\n\n'
				text.text += txt

			# Комментарий:
			m = txt.startswith(comment)
			if m:
				slide = prs.slides.add_slide(text_slide_layout)
				title = slide.shapes.title
				text = slide.placeholders[1]

				title.text = 'Комментарий'
				text.text = txt.partition(':')[2]

	# save to memory
	if file_stream == None:
		file_stream = io.BytesIO()
	prs.save(file_stream)
	return file_stream

# save pptx to file
def save_pptx(file_stream, pack_title):
	file_name = unidecode(pack_title)
	file_name = re.sub(r'[^\w\-_\. ]', '_', file_name)
	file_name += '.pptx'
	prs = Presentation(file_stream)
	prs.save(file_name)
	print ('-> ', file_name)

# main()
if len(sys.argv) != 2:
	print('Usage: ', sys.argv[0], ' url')
	sys.exit(0)

url = sys.argv[1]
urlo = urlparse(url)
r = requests.get(url, verify=False)
page = r.content.decode('utf-8')
soup = BeautifulSoup(page, 'lxml')
pack_title = soup.title.text

# если надо скачать все туры (only for kand.info)
if urlo.path.startswith('/tour/') and urlo.netloc == 'kand.info':
	file_stream = io.BytesIO()
	first_time = True
	for a in soup.find_all('a'):
		url_tour = a.attrs['href']
		if re.match(r'/node/\d+/tour/\d+', url_tour):
			url_tour = urljoin(url, url_tour)
			if first_time:
				file_stream = generate_pptx(url_tour, None)
				first_time = False
			else:
				file_stream = generate_pptx(url_tour, file_stream)
else:
	file_stream = generate_pptx(url)

save_pptx(file_stream, pack_title)